import fitz  # PyMuPDF
import docx
import email
import os
from typing import Optional
import pandas as pd
from pptx import Presentation
from app.services.rag import read_image


def sanitize_text(text: str) -> str:
    # Remove null characters and strip
    return text.replace("\x00", "").strip()

def extract_from_pdf(file_path: str) -> str:
    try:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return sanitize_text(text)
    except Exception as e:
        raise RuntimeError(f"PDF extraction failed: {e}")

def extract_from_docx(file_path: str) -> str:
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return sanitize_text(text)
    except Exception as e:
        raise RuntimeError(f"DOCX extraction failed: {e}")

def extract_from_email(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            msg = email.message_from_file(f)
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            body += payload.decode(errors="ignore")
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    body += payload.decode(errors="ignore")
            return sanitize_text(body)
    except Exception as e:
        raise RuntimeError(f"Email extraction failed: {e}")

def extract_from_csv(file_path: str) -> str:
    try:
        df = pd.read_csv(file_path, encoding="utf-8", errors="ignore")
        return sanitize_text(df.to_markdown(index=False))
    except Exception as e:
        raise RuntimeError(f"CSV extraction failed: {e}")

def extract_from_xlsx(file_path: str) -> str:
    try:
        df_list = pd.read_excel(file_path, sheet_name=None) 
        text = ""
        for sheet_name, df in df_list.items():
            text += f"Sheet: {sheet_name}\n"
            text += df.to_markdown(index=False)
            text += "\n"
        return sanitize_text(text)
    except Exception as e:
        raise RuntimeError(f"XLSX extraction failed: {e}")      

EXT_TO_MIME = {
    "jpeg": "image/jpeg",
    "jpg": "image/jpeg",
    "png": "image/png",
    "gif": "image/gif",
    "bmp": "image/bmp",
    "tiff": "image/tiff",
    "svg": "image/svg+xml",
}

def extract_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    extracted_content = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = {
            'slide_number': slide_num,
            'text': '',
            'images': []
        }

        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_parts.append(shape.text)
        
        slide_content['text'] = '\n'.join(text_parts)
        
        for shape in slide.shapes:
            if shape.shape_type == 13: 
                image = shape.image
                image_bytes = image.blob
                ext = image.ext.lower()
                mime_type = EXT_TO_MIME.get(ext, "application/octet-stream")
                image_content = read_image(image_bytes=image_bytes, mime_type=mime_type)
                
                slide_content['images'].append(image_content)
        
        extracted_content.append(slide_content)
    
    output_lines = []

    for slide in extracted_content:
        slide_str = f"Slide Number: {slide['slide_number']}\t"
        slide_str += f"Text: {slide['text']}\t"
        
        for idx, image in enumerate(slide['images'], 1):
            slide_str += f"Related_image_description_{idx}: {image}\t"

        output_lines.append(slide_str.strip())

    return "\n\n".join(output_lines)

def extract_fallback(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return sanitize_text(f.read())
    except Exception:
        try:
            with open(file_path, "rb") as f:
                content = f.read()
                return sanitize_text(content.decode("utf-8", errors="ignore"))
        except Exception as e:
            raise RuntimeError(f"Fallback extraction failed: {e}")

def extract_text(file_path: str, mime_type: Optional[str] = None) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            return extract_from_pdf(file_path)
        elif ext == ".docx" or mime_type in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ]:
            return extract_from_docx(file_path)
        elif ext in [".eml", ".msg"] or mime_type == "message/rfc822":
            return extract_from_email(file_path)
        elif ext == ".csv" or mime_type == "text/csv":
            return extract_from_csv(file_path)
        elif ext == ".xlsx" or mime_type in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ]:
            return extract_from_xlsx(file_path)
        elif ext == ".pptx" or mime_type in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            return extract_from_pptx(file_path)
        else:
            return extract_fallback(file_path)
    except Exception as e:
        return f"Error extracting text: {str(e)}"